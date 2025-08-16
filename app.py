from flask import Flask, request, render_template, send_file
from pptx import Presentation
import pandas as pd
import re
import os
from datetime import datetime

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# --- Regex for dynamic KPI extraction ---
KPI_FINDER = re.compile(r"""
    (?P<kpi>[A-Za-z][A-Za-z0-9\s/&\-()]+?)      # KPI name (letters, spaces, symbols)
    (?:\s*[:\-–—]\s*|\s+is\s+|\s+was\s+|\s+)    # separator: :, -, is, was, or space
    (?P<value>
        \(?\s*                                   # optional opening (
        (?:[₹$€£]\s*)?                           # optional currency symbol
        [\d,]+(?:\.\d+)?                         # numbers with , and . support
        (?:\s?(?:cr|crore|lakh|m|mn|b|bn))?      # optional magnitude
        %?                                       # optional %
        \s*\)?                                   # optional closing )
    )
""", re.IGNORECASE | re.VERBOSE)

# --- Helper to merge runs inside shapes ---
def extract_runs_text(shape):
    text_parts = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text_parts.append(run.text.strip())
    return " ".join(text_parts).strip()

# --- Extract KPIs from one PPT ---
def extract_kpis(ppt_file, filename):
    prs = Presentation(ppt_file)
    rows = []
    manager = os.path.splitext(filename)[0]

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            text = extract_runs_text(shape)
            if not text:
                continue
            for m in KPI_FINDER.finditer(text):
                kpi = m.group("kpi").strip().replace("\n", " ")
                value = m.group("value").strip()
                rows.append({
                    "Manager/Sector": manager,
                    "Slide": slide_idx,
                    "KPI": kpi,
                    "Value": value
                })

    if rows:
        df = pd.DataFrame(rows).drop_duplicates(
            subset=["Manager/Sector", "Slide", "KPI", "Value"]
        )
        return df
    return pd.DataFrame(columns=["Manager/Sector", "Slide", "KPI", "Value"])

# --- Flask route ---
@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        files = request.files.getlist("pptfiles")
        all_data = []

        for ppt_file in files:
            if ppt_file and ppt_file.filename.endswith(".pptx"):
                filename = ppt_file.filename
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                ppt_file.save(file_path)
                df = extract_kpis(file_path, filename)
                all_data.append(df)

        if all_data:
            final_df = pd.concat(all_data, ignore_index=True)

            # 👉 Pivot: one row per Manager, KPI names as columns
            pivot_df = final_df.pivot_table(
                index="Manager/Sector",
                columns="KPI",
                values="Value",
                aggfunc="first"
            ).reset_index()

            # Save to Excel
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = os.path.join(
                UPLOAD_FOLDER, f"KPI_Pivot_{timestamp}.xlsx"
            )
            pivot_df.to_excel(output_file, index=False)

            return send_file(output_file, as_attachment=True)

    return render_template("index.html")

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000)
